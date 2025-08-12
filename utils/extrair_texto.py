import pandas as pd
import fitz  # PyMuPDF
import docx
from pptx import Presentation
from io import BytesIO
import streamlit as st

@st.cache_data(show_spinner="Extraindo texto...")
def extrair_texto(file, filename: str):
    """
    Função principal para extrair texto de vários formatos de arquivo.
    Detecta a extensão do arquivo e chama a função específica.
    """
    ext = filename.split('.')[-1].lower()

    if ext == "pdf":
        return extract_pdf(file)
    elif ext == "docx":
        return extract_docx(file)
    elif ext == "txt":
        return file.read().decode("utf-8")
    elif ext == "pptx":
        return extract_pptx(file)
    elif ext == "xlsx":
        return extract_xlsx(file)
    else:
        return "Formato não suportado."

@st.cache_data
def extract_pdf(file_obj):
    """
    Extrai texto de arquivo PDF usando PyMuPDF (fitz).
    Lê todas as páginas e concatena o texto.
    """
    file_bytes = file_obj.read()
    with fitz.open(stream=file_bytes, filetype="pdf") as doc:
        return "".join([page.get_text() for page in doc])

@st.cache_data
def extract_docx(file_obj):
    """
    Extrai texto de arquivo DOCX usando python-docx.
    Junta os textos de todos os parágrafos e tabelas.
    """
    doc = docx.Document(file_obj)
    textos = []

    # Texto dos parágrafos
    for p in doc.paragraphs:
        if p.text.strip():
            textos.append(p.text.strip())

    # Texto das tabelas
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    textos.append(cell.text.strip())

    return "\n".join(textos)

@st.cache_data
def extract_pptx(file_obj):
    """
    Extrai texto de apresentações PPTX usando python-pptx.
    Percorre slides e formas, extraindo texto de cada forma que tenha atributo text.
    """
    texto = ""
    for slide in Presentation(file_obj).slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texto += shape.text + "\n"
    return texto

@st.cache_data
def extract_xlsx(file_obj):
    """
    Extrai texto de planilhas XLSX usando pandas.
    Para cada aba, converte o conteúdo em string e concatena com cabeçalho da planilha.
    """
    xls = pd.ExcelFile(file_obj)
    texto = ""
    for sheet in xls.sheet_names:
        df = xls.parse(sheet)
        texto += f"\n--- Planilha: {sheet} ---\n"
        texto += df.astype(str).to_string(index=False)
    return texto
